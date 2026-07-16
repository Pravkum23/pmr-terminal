"""PMR Terminal - daily PowerPoint deck with sparklines + risk scatter."""
from __future__ import annotations

import io

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Inches, Pt

BG = RGBColor(0x0A, 0x0E, 0x1A)
PANEL = RGBColor(0x11, 0x18, 0x27)
TXT = RGBColor(0xD7, 0xE1, 0xF3)
DIM = RGBColor(0x7C, 0x8D, 0xB5)
GREEN = RGBColor(0x22, 0xC5, 0x5E)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
RED = RGBColor(0xEF, 0x44, 0x44)
CYAN = RGBColor(0x22, 0xD3, 0xEE)
RAGC = {"GREEN": GREEN, "AMBER": AMBER, "RED": RED}


def _slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    return s


def _text(s, x, y, w, h, txt, size=12, color=TXT, bold=False, align=None):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = txt
    r.font.size, r.font.bold, r.font.color.rgb = Pt(size), bold, color
    r.font.name = "Segoe UI"
    if align:
        from pptx.enum.text import PP_ALIGN
        p.alignment = getattr(PP_ALIGN, align)
    return tb


def _title(s, main, sub):
    _text(s, 0.4, 0.18, 9, 0.5, main, 22, CYAN, True)
    _text(s, 0.4, 0.68, 12, 0.35, sub, 11, DIM)


def _sparkline_png(vals, up: bool) -> io.BytesIO:
    fig, ax = plt.subplots(figsize=(1.7, 0.42), dpi=110)
    ax.plot(vals, color="#22c55e" if up else "#ef4444", lw=1.3)
    ax.axis("off"); fig.patch.set_alpha(0)
    buf = io.BytesIO()
    fig.savefig(buf, format="png", transparent=True, bbox_inches="tight", pad_inches=0.02)
    plt.close(fig); buf.seek(0)
    return buf


def _pct_run(p, x, suffix="%"):
    r = p.add_run()
    if x is None:
        r.text = "–"; r.font.color.rgb = DIM
    else:
        r.text = f"{x:+.1f}{suffix}"
        r.font.color.rgb = GREEN if x >= 0 else RED
    r.font.size = Pt(10); r.font.name = "Segoe UI"
    return r


def build_deck(d: dict, path: str = "docs/pmr_daily.pptx"):
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.33), Inches(7.5)
    b, m, p = d["breadth"], d["movers"], d["pnl"]
    cur = "₹" if p["base_currency"] == "INR" else "$"

    # ---- Slide 1: cover / summary ----
    s = _slide(prs)
    _text(s, 0.6, 1.7, 12, 1.0, "PMR TERMINAL", 48, CYAN, True)
    _text(s, 0.62, 2.7, 12, 0.5, "Daily Market Intelligence · Praveen · Moulya · Risha", 16, TXT)
    _text(s, 0.62, 3.3, 12, 0.4, d["generated_at"], 12, DIM)
    tone_c = GREEN if b["tone"] == "RISK-ON" else RED if b["tone"] == "RISK-OFF" else AMBER
    _text(s, 0.62, 4.2, 6, 0.6, f"Market tone: {b['tone']}", 20, tone_c, True)
    _text(s, 0.62, 4.9, 12, 0.5,
          f"RAG {b['red']}R / {b['amber']}A / {b['green']}G   ·   "
          f"{b['pct_above_200dma']}% above 200DMA   ·   "
          f"Portfolio {cur}{p['total_value']:,.0f} ({p['total_pnl_pct']:+.1f}%)", 13, TXT)
    _text(s, 0.62, 6.9, 12, 0.4, "Educational analytics — not investment advice.", 9, DIM)

    # ---- Slide 2: MTD movers with sparklines ----
    s = _slide(prs)
    _title(s, "MTD Gainers & Losers", "Top 5 each · 3-month sparkline")
    for col, (label, items) in enumerate((("GAINERS", m["gainers"]), ("LOSERS", m["losers"]))):
        x0 = 0.5 + col * 6.5
        _text(s, x0, 1.15, 3, 0.35, label, 13, GREEN if col == 0 else RED, True)
        for i, r in enumerate(items):
            y = 1.65 + i * 1.05
            _text(s, x0, y, 3.2, 0.35, r["name"], 13, TXT, True)
            _text(s, x0, y + 0.33, 3.2, 0.3, f"{r['symbol']} · {r['rag']}", 9,
                  RAGC[r["rag"]])
            tb = s.shapes.add_textbox(Inches(x0 + 3.3), Inches(y + 0.05), Inches(1.2), Inches(0.4))
            _pct_run(tb.text_frame.paragraphs[0], r["mtd"])
            if r.get("spark"):
                img = _sparkline_png(r["spark"], (r["mtd"] or 0) >= 0)
                s.shapes.add_picture(img, Inches(x0 + 4.6), Inches(y), height=Inches(0.42))

    # ---- Slide 3: risk scatter ----
    s = _slide(prs)
    _title(s, "Risk Map", "1Y return vs annualized volatility · color = RAG · size = |max drawdown|")
    pts = [r for r in d["risk"] if r.get("ret_1y") is not None and r.get("vol_ann")]
    fig, ax = plt.subplots(figsize=(11.5, 5.4), dpi=110)
    fig.patch.set_facecolor("#0a0e1a"); ax.set_facecolor("#0d1424")
    cmap = {"GREEN": "#22c55e", "AMBER": "#f59e0b", "RED": "#ef4444"}
    ax.scatter([r["vol_ann"] for r in pts], [r["ret_1y"] for r in pts],
               s=[max(30, abs(r["max_dd_1y"]) * 8) for r in pts],
               c=[cmap[r["rag"]] for r in pts], alpha=0.65, edgecolors="white", linewidths=0.4)
    for r in pts:
        ax.annotate(r["symbol"], (r["vol_ann"], r["ret_1y"]), fontsize=6.5,
                    color="#7c8db5", xytext=(4, 3), textcoords="offset points")
    ax.axhline(0, color="#1f2a44", lw=1)
    ax.set_xlabel("Annualized Vol %", color="#7c8db5"); ax.set_ylabel("1Y Return %", color="#7c8db5")
    ax.tick_params(colors="#7c8db5"); [sp.set_color("#1f2a44") for sp in ax.spines.values()]
    buf = io.BytesIO(); fig.savefig(buf, format="png", facecolor="#0a0e1a",
                                    bbox_inches="tight"); plt.close(fig); buf.seek(0)
    s.shapes.add_picture(buf, Inches(0.5), Inches(1.3), width=Inches(12.3))

    # ---- Slide 4: AI signals ----
    s = _slide(prs)
    _title(s, "AI Signal Desk", "Buy/Sell engine — top conviction calls")
    actionable = [r for r in d["signals"] if r["signal"] != "HOLD"][:12] or d["signals"][:12]
    for i, r in enumerate(actionable):
        y = 1.2 + i * 0.47
        _text(s, 0.5, y, 3.4, 0.35, r["name"], 12, TXT, True)
        sig_c = GREEN if "BUY" in r["signal"] else RED
        _text(s, 3.9, y, 1.8, 0.35, r["signal"], 12, sig_c, True)
        _text(s, 5.7, y, 1.2, 0.35, f"{int(r['confidence'])}%", 11, DIM)
        _text(s, 6.9, y, 6.2, 0.35, " · ".join(r["signal_reasons"][:3]), 9, DIM)

    # ---- Slide 5: portfolio ----
    s = _slide(prs)
    _title(s, "Portfolio P&L", f"Base currency {p['base_currency']} · USD/INR {p['usdinr']}")
    hdr = ["Symbol", "Qty", "Avg", "Last", "Value", "P&L", "P&L %", "Day %"]
    xs = [0.5, 2.6, 3.6, 4.7, 5.9, 7.6, 9.3, 10.6]
    for x, h in zip(xs, hdr):
        _text(s, x, 1.15, 1.5, 0.3, h, 10, DIM, True)
    for i, x0 in enumerate(p["positions"][:14]):
        y = 1.5 + i * 0.38
        vals = [x0["symbol"], f"{x0['qty']}", f"{x0['avg_price']:,.0f}", f"{x0['price']:,.0f}",
                f"{cur}{x0['value']:,.0f}", f"{cur}{x0['pnl']:,.0f}",
                f"{x0['pnl_pct']:+.1f}%", f"{x0['day_pct']:+.1f}%"]
        for j, (x, v) in enumerate(zip(xs, vals)):
            c = TXT
            if j == 5 or j == 6:
                c = GREEN if x0["pnl"] >= 0 else RED
            if j == 7:
                c = GREEN if x0["day_pct"] >= 0 else RED
            _text(s, x, y, 1.9, 0.3, v, 10, c, j == 0)
    rk = p.get("risk") or {}
    _text(s, 0.5, 7.0, 12, 0.35,
          f"Total {cur}{p['total_value']:,.0f} · P&L {cur}{p['total_pnl']:,.0f} ({p['total_pnl_pct']:+.1f}%) · "
          f"Sharpe {rk.get('sharpe','–')} · Vol {rk.get('vol_ann','–')}% · VaR95 {rk.get('var95_daily','–')}%/day",
          11, CYAN, True)

    prs.save(path)
    print(f"Deck saved: {path}")
